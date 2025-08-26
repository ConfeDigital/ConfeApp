"""
Proyecto de Vida report generator.
Handles the complete generation of life project PowerPoint presentations.
"""
from io import BytesIO
from django.http import HttpResponse
from datetime import datetime
from candidatos.models import UserProfile
from .data_collector import ReportDataCollector
import json

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


class ProyectoVidaReport:
    """Generator for Proyecto de Vida PowerPoint presentations."""
    
    def __init__(self):
        self.navy_color = RGBColor(6, 45, 85)  # #062d55
        self.white_color = RGBColor(255, 255, 255)
        self.light_blue = RGBColor(173, 216, 230)  # Light blue for accents
    
    def get_profile(self, uid):
        """Get user profile."""
        try:
            return UserProfile.objects.get(user__id=uid)
        except UserProfile.DoesNotExist:
            raise ValueError(f"Profile not found for user ID: {uid}")
    
    def create_title_slide(self, prs, profile):
        """Create the title slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.navy_color
        
        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "PROYECTO DE VIDA"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = self.white_color
        
        # Subtitle
        full_name = f"{profile.user.first_name} {profile.user.last_name}"
        if hasattr(profile.user, 'second_last_name') and profile.user.second_last_name:
            full_name += f" {profile.user.second_last_name}"
        date_str = datetime.now().strftime("%d/%m/%Y")
        
        sub_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.text = f"Nombre: {full_name}\nFecha: {date_str}"
        for p in sub_tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = self.white_color
    
    def create_section_slide(self, prs, title, content):
        """Create a section slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        # Header bar
        header = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.navy_color
        header.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.white_color
        
        # Content
        content_box = shapes.add_textbox(Inches(1), Inches(1.3), prs.slide_width - Inches(2), prs.slide_height - Inches(2))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_top = Inches(0.1)
        tf_content.margin_left = Inches(0.1)
        tf_content.margin_right = Inches(0.1)
        tf_content.text = content
        
        for paragraph in tf_content.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 0, 0)

    def create_meta_slide(self, prs, title, meta_data):
        """Create a specially formatted slide for Meta sections."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes
        
        # Header bar
        header = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = self.navy_color
        header.line.fill.background()
        
        # Title
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = self.white_color
        
        # Meta description box with background
        meta_box = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), prs.slide_width - Inches(2), Inches(1))
        meta_fill = meta_box.fill
        meta_fill.solid()
        meta_fill.fore_color.rgb = self.light_blue
        meta_box.line.color.rgb = self.navy_color
        meta_box.line.width = Pt(2)
        
        # Meta text
        meta_text_box = shapes.add_textbox(Inches(1.2), Inches(1.7), prs.slide_width - Inches(2.4), Inches(0.6))
        meta_tf = meta_text_box.text_frame
        meta_tf.text = f"META: {meta_data.get('meta', 'No especificado')}"
        for paragraph in meta_tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(20)
                run.font.bold = True
                run.font.color.rgb = self.navy_color
        
        # Steps header
        steps_header_box = shapes.add_textbox(Inches(1), Inches(2.8), prs.slide_width - Inches(2), Inches(0.4))
        steps_header_tf = steps_header_box.text_frame
        steps_header_tf.text = "PASOS PARA ALCANZAR LA META:"
        for paragraph in steps_header_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = self.navy_color
        
        # Steps content
        pasos = meta_data.get("pasos", [])
        if pasos:
            y_position = 3.3
            for i, paso in enumerate(pasos, 1):
                # Step number circle
                circle = shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y_position), Inches(0.4), Inches(0.4))
                circle_fill = circle.fill
                circle_fill.solid()
                circle_fill.fore_color.rgb = self.navy_color
                circle.line.fill.background()
                
                # Step number text
                num_box = shapes.add_textbox(Inches(1.2), Inches(y_position), Inches(0.4), Inches(0.4))
                num_tf = num_box.text_frame
                num_tf.text = str(i)
                for paragraph in num_tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.bold = True
                        run.font.color.rgb = self.white_color
                
                # Step description
                desc_box = shapes.add_textbox(Inches(1.8), Inches(y_position), Inches(4), Inches(0.4))
                desc_tf = desc_box.text_frame
                desc_tf.text = paso.get("descripcion", 'Sin descripción')
                for paragraph in desc_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                
                # Responsible person
                resp_box = shapes.add_textbox(Inches(6), Inches(y_position), Inches(2.5), Inches(0.4))
                resp_tf = resp_box.text_frame
                resp_tf.text = f"Encargado: {paso.get('encargado', 'No especificado')}"
                for paragraph in resp_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)
                        run.font.italic = True
                        run.font.color.rgb = RGBColor(100, 100, 100)
                
                y_position += 0.7
        else:
            # No steps message
            no_steps_box = shapes.add_textbox(Inches(1), Inches(3.3), prs.slide_width - Inches(2), Inches(0.5))
            no_steps_tf = no_steps_box.text_frame
            no_steps_tf.text = "No se han definido pasos para esta meta."
            for paragraph in no_steps_tf.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(150, 150, 150)
    
    def format_meta_content(self, meta_data):
        """Format meta content for display."""
        try:
            if isinstance(meta_data, str):
                # Try JSON first
                try:
                    return json.loads(meta_data)
                except json.JSONDecodeError:
                    # Try Python literal
                    try:
                        return ast.literal_eval(meta_data)
                    except Exception:
                        return {"meta": meta_data, "pasos": []}
            elif isinstance(meta_data, dict):
                # If the 'meta' key itself is a stringified dict
                if isinstance(meta_data.get("meta"), str) and meta_data["meta"].startswith("{"):
                    try:
                        return ast.literal_eval(meta_data["meta"])
                    except Exception:
                        return meta_data
                return meta_data
            else:
                return {"meta": str(meta_data), "pasos": []}
        except Exception:
            return {"meta": "No especificado", "pasos": []}
    
    def get_sections_data(self, data_collector):
        """Get all sections data for the presentation."""
        proyecto_vida_data = data_collector.get_proyecto_vida_data()
        
        sections = {
            "Grupo de Apoyo": proyecto_vida_data.get("Mi grupo de apoyo", "No especificado"),
            "Mis Talentos": proyecto_vida_data.get("Mis talentos personales", "No especificado"),
            "Talentos de Mi Grupo de Apoyo": proyecto_vida_data.get("Las cosas que a otras personas les gustan de mi", "No especificado"),
            "Lo Más Importante Para Mí": proyecto_vida_data.get("Lo más importante para mi", "No especificado"),
            "Cosas Sobre Mí": proyecto_vida_data.get("Cosas sobre mí", "No especificado"),
            "Apoyos Que Necesito": proyecto_vida_data.get("Lo que los demás deben de saber para apoyarme", "No especificado"),
            "Mi Historia": proyecto_vida_data.get("Mi historia, recuerdos que quiero compartir", "No especificado"),
            "Futuro Ideal": proyecto_vida_data.get("¿Cuál sería tu futuro ideal?", "No especificado"),
        }
        
        # Special handling for Meta sections
        meta_sections = {
            "Meta 1": proyecto_vida_data.get("Meta 1", {}),
            "Meta 2": proyecto_vida_data.get("Meta 2", {}),
            "Meta 3": proyecto_vida_data.get("Meta 3", {}),
        }
        
        return sections, meta_sections
    
    def generate(self, uid):
        """Generate the complete Proyecto de Vida presentation."""
        # Get profile and data collector
        profile = self.get_profile(uid)
        data_collector = ReportDataCollector(uid)
        
        # Create presentation
        prs = Presentation()
        
        # Create title slide
        self.create_title_slide(prs, profile)
        
        # Get sections data
        sections, meta_sections = self.get_sections_data(data_collector)
        
        # Create regular section slides
        for title, content in sections.items():
            self.create_section_slide(prs, title, content)
        
        # Create special Meta slides
        for title, meta_data in meta_sections.items():
            formatted_meta = self.format_meta_content(meta_data)

            # If meta itself is a string with dict inside, parse again
            if isinstance(formatted_meta.get("meta"), str) and formatted_meta.get("meta").startswith("{"):
                try:
                    import ast
                    parsed = ast.literal_eval(formatted_meta["meta"])
                    formatted_meta = parsed
                except Exception:
                    pass

            self.create_meta_slide(prs, title, formatted_meta)

        
        # Export PowerPoint
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Create response
        response = HttpResponse(
            ppt_buffer.getvalue(), 
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = f'attachment; filename="proyecto_vida_{uid}.pptx"'
        
        return response