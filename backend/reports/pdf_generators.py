# pdf_generators.py

from io import BytesIO
from django.http import HttpResponse
from reportlab.lib.pagesizes import letter, landscape, A3
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, PageBreak, Frame
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from datetime import datetime
from reportlab.pdfgen import canvas 
from django.core.files.storage import default_storage
import os
import json
import requests
import re 
from .reports_utils import draw_logo_header, draw_logo_header_canvas, draw_footer, draw_table, redondear_percentil, get_habilidades_adaptativas_coloreadas
from .table_templates import TABLE_TEMPLATES
from .data_fillers import get_pv_answers, fill_table_data, map_answers_to_template
from cuestionarios.models import Respuesta, Cuestionario
from discapacidad.models import CHItem
import unicodedata
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph,
    Spacer, PageBreak, Frame, PageTemplate
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch
from io import BytesIO
import unicodedata
import re
from django.http import HttpResponse
import os
from django.conf import settings
from reportlab.lib.units import cm
from cuestionarios.utils import get_resumen_cuestionarios_completo


def generate_ficha_tecnica(uid, profile, respuestas):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, leftMargin=50, rightMargin=50, topMargin=40, bottomMargin=40)
    elements = []
    styles = getSampleStyleSheet()

    # === Estilos reutilizables ===
    title_style = ParagraphStyle("TitleStyle", fontSize=14, textColor=colors.black, fontName="Helvetica-Bold", alignment=0, spaceAfter=12)
    normal_style = ParagraphStyle("NormalStyle", fontSize=10, textColor=colors.black, fontName="Helvetica")

    def section_header(text):
        return Paragraph(f"<b><font size=14 color='white'>{text}</font></b>", ParagraphStyle(
            name="SectionTitle",
            backColor=colors.Color(6 / 255, 45 / 255, 85 / 255),
            alignment=1,
            spaceBefore=6,
            spaceAfter=12,
            leading=16,
            textColor=colors.white
        ))
    estilo_tabla = TableStyle([
        ('GRID', (0, 0), (-1, -1), 0.25, colors.grey),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
    ])
    # === PÁGINA 1 ===
    elements.append(draw_logo_header())
    elements.append(Spacer(1, 10))

    # Título y foto
    profile_pic = None
    if profile.photo and default_storage.exists(profile.photo.name):
        img_path = default_storage.path(profile.photo.name)
        profile_pic = Image(img_path, width=100, height=100, kind='proportional')

    header_table_data = [[
        Table([
            [Paragraph("Institución CONFE a Favor de la Persona con Discapacidad Intelectual I.A.P.", title_style)],
            [Paragraph("CENTRO ESPECIALIZADO DE INCLUSIÓN LABORAL", title_style)],
            [Paragraph("FICHA TÉCNICA", title_style)]
        ], colWidths=[400]),
        profile_pic if profile_pic else ""
    ]]
    header_table = Table(header_table_data, colWidths=[400, 100])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('ALIGN', (-1, 0), (-1, -1), 'RIGHT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
    ]))
    elements.append(header_table)
    elements.append(Spacer(1, 12))

    elements.append(section_header("DATOS PERSONALES"))

    discapacidades = ", ".join([d.name for d in profile.disability.all()]) if profile.disability.exists() else "No especificada"
    full_name = f"{profile.user.first_name} {profile.user.last_name} {getattr(profile.user, 'second_last_name', '')}"
    gender_display = dict(profile.GENDER_CHOICES).get(profile.gender, "No especificado")
    stage_display = dict(profile.STAGE_CHOICES).get(profile.stage, "No especificado")

    info_data = [
        ["Nombre completo:", full_name, "Género:", gender_display],
        ["Discapacidad:", discapacidades, "Etapa:", stage_display],
        ["Teléfono:", profile.phone_number or "N/A", "", ""],
        ["Edad:", str(datetime.now().year - profile.birth_date.year) + " años" if profile.birth_date else "N/A",
         "Tipo de sangre:", profile.blood_type or "N/A"],
        ["CURP:", profile.curp or "N/A",
         "Medicamentos:", ", ".join([m.name for m in profile.medications.all()]) if profile.medications.exists() else "N/A"],
        ["Registro:", profile.registration_date.strftime('%Y-%m-%d') if profile.registration_date else "N/A",
         "Alergias:", profile.allergies or "N/A"],
        ["Presenta convulsiones:", "Sí" if profile.has_seizures else "No",
         "Recibe pensión:", "Sí" if profile.receives_pension else "No"],
        ["Tipo de pensión:", profile.pension_type if hasattr(profile, 'pension_type') and profile.pension_type else "N/A", "", ""],
        ["Cert. de discapacidad:", "Sí" if profile.has_disability_certificate else "No",
         "Juicio de interdicción:", "Sí" if profile.has_interdiction_judgment else "No"],
        ["Atención psicológica/\npsiquiátrica:", f"{'Sí' if profile.receives_psychological_care else 'No'} / {'Sí' if profile.receives_psychiatric_care else 'No'}", "", ""],
        ["Dirección:", (
            f"{profile.domicile.address_road}, {profile.domicile.address_number}, "
            f"{profile.domicile.address_col}, {profile.domicile.address_municip}, "
            f"{profile.domicile.address_city}, {profile.domicile.address_state}, CP {profile.domicile.address_PC}"
        ) if profile.domicile else "N/A", "", ""]
    ]
    info_table = Table(info_data, colWidths=[125, 235, 110, 100])
    info_table.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (2, 0), (2, -1), 'Helvetica-Bold'),
        ('BACKGROUND', (0, 0), (-1, -1), colors.whitesmoke),
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('LEFTPADDING', (0, 0), (-1, -1), 4),
        ('RIGHTPADDING', (0, 0), (-1, -1), 4),
    ]))
    elements.append(info_table)
    elements.append(Spacer(1, 10))

    elements.append(section_header("CONTACTO(S)"))
    if profile.emergency_contacts.exists():
        for contact in profile.emergency_contacts.all():
            name = f"{contact.first_name} {contact.last_name} {contact.second_last_name or ''}"
            relationship = dict(contact.RELATIONSHIP_CHOICES).get(contact.relationship, "N/A")
            address = contact.domicile.__str__() if contact.domicile else "N/A"
            lives_same = "Sí" if contact.lives_at_same_address else "No"
            contact_str = f"{name} ({relationship}) - Tel: {contact.phone_number} - Mismo domicilio: {lives_same} - Dirección: {address}"
            elements.append(Paragraph(contact_str, normal_style))
    else:
        elements.append(Paragraph("N/A", normal_style))
    elements.append(Spacer(1, 12))

    tables = fill_table_data(profile, TABLE_TEMPLATES)

    resumen_completo = get_resumen_cuestionarios_completo(profile.user.id)
    datos_proyecto_vida = resumen_completo.get("proyecto_vida", {})
    datos_entrevista = resumen_completo.get("entrevista", {})
    datos_diagnostica = resumen_completo.get("evaluacion_diagnostica", {})

    elements.append(section_header("EVALUACIÓN DIAGNÓSTICA"))
    col1, col2 = [], []
    for i, (pregunta, respuesta) in enumerate(datos_diagnostica.items()):
        p = Paragraph(f"<b>{pregunta}</b>: {respuesta}", normal_style)
        (col1 if i % 2 == 0 else col2).append(p)
        (col1 if i % 2 == 0 else col2).append(Spacer(1, 4))
    diag_columns = Table([[col1, col2]], colWidths=[250, 250])
    diag_columns.setStyle(TableStyle([
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('TOPPADDING', (0, 0), (-1, -1), 0),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 0),
    ]))
    elements.append(diag_columns)
    elements.append(PageBreak())

    for section in ["NECESIDADES DE APOYO", "PROYECTO DE VIDA", "TALENTOS"]:
        elements.append(section_header(section))
        elements.append(Spacer(1, 6))

        if section == "NECESIDADES DE APOYO":
            entrevista_data = [
                ["¿Cómo te ves en tu futuro? ¿Qué metas te gustaría cumplir?", datos_entrevista.get("futuro_usuario") or ""],
                ["¿A futuro cómo le gustaría ver a su hijo/hija?", datos_entrevista.get("futuro_hijo") or ""],
                ["Observaciones del entrevistador", datos_entrevista.get("observaciones_entrevistador") or ""],
            ]
            data = [[Paragraph(preg, normal_style), Paragraph(resp, normal_style)] for preg, resp in entrevista_data]
            table = Table(data, colWidths=[200, 250])
            table.setStyle(estilo_tabla)
            elements.append(table)

        elif section == "PROYECTO DE VIDA":
            textos_data = [
                ["Lo más importante para mí", datos_proyecto_vida.get("lo_mas_importante") or ""],
                ["Me gusta, me tranquiliza, me hace sentir bien, me divierte", datos_proyecto_vida.get("me_gusta") or ""],
            ]
            data = [[Paragraph(preg, normal_style), Paragraph(resp, normal_style)] for preg, resp in textos_data]
            table = Table(data, colWidths=[200, 250])
            table.setStyle(estilo_tabla)
            elements.append(table)
            elements.append(Spacer(1, 10))

            if datos_proyecto_vida["metas"]:
                elements.append(Paragraph("Metas", title_style))
                metas_data = []
                for meta in datos_proyecto_vida["metas"]:
                    pasos_render = "<br/>".join(
                        f"- {p['descripcion']} <i>({p['encargado']})</i>" for p in meta["pasos"]
                    )
                    texto_meta = f"<b>Meta:</b> {meta['meta']}<br/><b>Pasos:</b><br/>{pasos_render}"
                    metas_data.append([Paragraph(texto_meta, normal_style)])
                table = Table(metas_data, colWidths=[450])
                table.setStyle(estilo_tabla)
                elements.append(table)
            elements.append(Spacer(1, 10))

        elif section == "TALENTOS":
            talentos = datos_proyecto_vida.get("talentos", {})
            if talentos:
                talentos_data = []
                for pregunta, respuestas in talentos.items():
                    if respuestas:
                        pregunta_parrafo = Paragraph(f"<b>{pregunta}</b>", normal_style)
                        respuesta_parrafo = Paragraph(", ".join(respuestas), normal_style)
                        talentos_data.append([pregunta_parrafo, respuesta_parrafo])
                if talentos_data:
                    table = Table(talentos_data, colWidths=[200, 250])
                    table.setStyle(estilo_tabla)
                    elements.append(table)
                else:
                    elements.append(Paragraph("No se registraron talentos con contenido en esta sección.", normal_style))
            else:
                elements.append(Paragraph("No se registraron talentos para este usuario.", normal_style))

        elements.append(Spacer(1, 20))

    # === Página 3: SECCIÓN SIS ===
    elements.append(PageBreak())
    # Título y logo SIS
    sis_title = Paragraph("<b><font size=16>ESCALA DE INTENSIDAD DE APOYOS (SIS)</font></b>", ParagraphStyle(
    name="SISTitle",
    alignment=0,
    spaceAfter=10,
    fontName="Helvetica-Bold",
    fontSize=16,
    ))

    logo_ceil_path = os.path.join(settings.MEDIA_ROOT, "logos/sis.png")
    logo_jap_path = os.path.join(settings.MEDIA_ROOT, "logos/sis2.png")

    logo_ceil = Image(logo_ceil_path, width=50, height=50) if os.path.exists(logo_ceil_path) else Paragraph("", normal_style)
    logo_jap = Image(logo_jap_path, width=50, height=50) if os.path.exists(logo_jap_path) else Paragraph("", normal_style)

    sis_header_row = [[sis_title, logo_ceil, logo_jap]]
    sis_header = Table(sis_header_row, colWidths=[260, 100, 100])
    sis_header.setStyle(TableStyle([
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('ALIGN', (0, 0), (0, 0), 'CENTER'),
        ('ALIGN', (1, 0), (1, 0), 'CENTER'),
        ('ALIGN', (2, 0), (2, 0), 'CENTER'),
    ]))
    elements.append(sis_header)
    elements.append(Spacer(1, 24))

    updated_table, celdas_coloreadas = get_habilidades_adaptativas_coloreadas(uid, tables["Habilidades Adaptativas - Tabla de resultados SIS"])
    elements.extend(draw_table(updated_table, "", celdas_coloreadas))
    elements.append(PageBreak())

    elements.append(section_header("PROTECCIÓN Y DEFENSA"))
    prot_table = draw_table(tables["Protección y Defensa"], "PROTECCIÓN Y DEFENSA")
    # Asegurar que se usen colWidths más amplios (no los de percentiles):
    for t in prot_table:
        if hasattr(t, 'colWidths'):
            t._argW = [10 * cm for _ in t._argW]  # Ancho amplio uniforme
    elements.extend(prot_table)
    elements.append(Spacer(1, 24))

    elements.append(section_header("NECESIDADES MÉDICAS Y CONDUCTUALES"))
    med_table = draw_table(tables["Necesidades Médicas y Conductuales"], "NECESIDADES MÉDICAS Y CONDUCTUALES")
    for t in med_table:
        if hasattr(t, 'colWidths'):
            t._argW = [10*cm for _ in t._argW]  # Ancho amplio uniforme
    elements.extend(med_table)
    elements.append(Spacer(1, 12))

    doc.build(elements)
    buffer.seek(0)
    response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="ficha_tecnica_{uid}.pdf"'
    return response

def generate_proyecto_vida_pdf(profile):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    def add_title_slide():
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(6, 45, 85)  # #062d55

        # Title
        title_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "PROYECTO DE VIDA"
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Subtitle
        full_name = f"{profile.user.first_name} {profile.user.last_name} {getattr(profile.user, 'second_last_name', '')}"
        date_str = datetime.now().strftime("%d/%m/%Y")

        sub_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
        sub_tf = sub_box.text_frame
        sub_tf.text = f"Nombre: {full_name}\nFecha: {date_str}"
        for p in sub_tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(255, 255, 255)

    def add_section_slide(title, content):
        slide = prs.slides.add_slide(blank_slide_layout)
        shapes = slide.shapes

        # Header bar
        header = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
        fill = header.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(6, 45, 85)
        header.line.fill.background()

        # Title
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Content
        content_box = shapes.add_textbox(Inches(1), Inches(1.3), prs.slide_width - Inches(2), prs.slide_height - Inches(2))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_top = Inches(0.1)
        tf_content.margin_left = Inches(0.1)
        tf_content.margin_right = Inches(0.1)
        tf_content.text = content

        for paragraph in tf_content.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Get sections
    answers_dict = get_pv_answers(profile)
    sections = {
        "Grupo de Apoyo": answers_dict.get("Mi grupo de apoyo", "No especificado"),
        "Mis Talentos": answers_dict.get("Mis talentos personales", "No especificado"),
        "Talentos de Mi Grupo de Apoyo": answers_dict.get("Las cosas que a otras personas les gustan de mi", "No especificado"),
        "Lo Más Importante Para Mí": answers_dict.get("Lo más importante para mi", "No especificado"),
        "Cosas Sobre Mí": answers_dict.get("Cosas sobre mí", "No especificado"),
        "Apoyos Que Necesito": answers_dict.get("Lo que los demás deben de saber para apoyarme", "No especificado"),
        "Mi historia": answers_dict.get("Mi historia, recuerdos que quiero compartir","No especificado"),
        "Meta 1": answers_dict.get("Meta 1", "No especificado"),
        "Meta 2": answers_dict.get("Meta 2", "No especificado"),
        "Meta 3": answers_dict.get("Meta 3", "No especificado"),
        "Futuro Ideal": answers_dict.get("¿Cuál sería tu futuro ideal?","No especificado"),
        "Pasos para Alcanzar Meta 1": answers_dict.get("Pasos para alcanzar Meta 1", "No especificado"),
        "Pasos para Alcanzar Meta 2": answers_dict.get("Pasos para alcanzar Meta 2", "No especificado"),
        "Pasos para Alcanzar Meta 3": answers_dict.get("Pasos para alcanzar Meta 3", "No especificado"),
    }

    # Generate slides
    add_title_slide()
    for title, content in sections.items():
        add_section_slide(title, content)

    # Export PowerPoint
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)

    response = HttpResponse(ppt_buffer.getvalue(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = f'attachment; filename="proyecto_vida_{profile.user.id}.pptx"'
    return response

def normalizar(texto):
    if not isinstance(texto, str):
        return ""
    texto = unicodedata.normalize("NFKD", texto)
    texto = texto.encode("ASCII", "ignore").decode("utf-8")
    texto = re.sub(r"\s+", " ", texto)
    return texto.strip().lower()


def normalizar(texto):
    if not isinstance(texto, str):
        return ""
    texto = unicodedata.normalize("NFKD", texto)
    texto = texto.encode("ASCII", "ignore").decode("utf-8")
    texto = re.sub(r"\s+", " ", texto)
    return texto.strip().lower()

def generate_cuadro_de_habilidades_pdf(profile):
    buffer = BytesIO()
    width, height = landscape(letter)
    styles = getSampleStyleSheet()
    NAVY = colors.Color(6 / 255, 45 / 255, 85 / 255)

    left_align = ParagraphStyle(name='LeftAlign', parent=styles['Normal'], alignment=0, fontSize=9)
    apoyo_style = ParagraphStyle(name='ApoyoStyle', parent=styles['Normal'], alignment=0, fontSize=8, leading=10)

    left_margin = 30
    right_margin = 30
    top_margin = 70
    bottom_margin = 40

    usable_width = width - left_margin - right_margin
    usable_height = height - top_margin - bottom_margin

    frame = Frame(left_margin, bottom_margin, usable_width, usable_height, id='normal')

    def draw_header(canvas, doc):
        canvas.setFillColor(NAVY)
        canvas.setFont("Helvetica-Bold", 16)
        canvas.drawCentredString(width / 2, height - 40, "CUADRO DE HABILIDADES")
        canvas.setFont("Helvetica", 12)
        canvas.drawCentredString(width / 2, height - 60,
                                 "Institución CONFE a Favor de la Persona con Discapacidad Intelectual I.A.P.")
        canvas.drawCentredString(width / 2, height - 75,
                                 "CENTRO ESPECIALIZADO DE INCLUSIÓN LABORAL")
        canvas.setStrokeColor(NAVY)
        canvas.line(30, height - 80, width - 30, height - 80)

    doc = SimpleDocTemplate(buffer, pagesize=landscape(letter),
                            leftMargin=left_margin, rightMargin=right_margin,
                            topMargin=top_margin, bottomMargin=bottom_margin)
    doc.addPageTemplates([PageTemplate(id='with-header', frames=frame, onPage=draw_header)])

    elements = []

    # Datos personales
    full_name = f"{profile.user.first_name} {profile.user.last_name} {getattr(profile.user, 'second_last_name', '')}"
    birth_date = profile.birth_date.strftime("%d/%m/%Y") if profile.birth_date else "N/A"
    registration_date = profile.registration_date.strftime("%Y-%m-%d") if profile.registration_date else "N/A"
    gender = dict(profile.GENDER_CHOICES).get(profile.gender, "No especificado")
    discapacidad = ", ".join([d.name for d in profile.disability.all()]) if profile.disability.exists() else "No especificada"
    stage = dict(profile.STAGE_CHOICES).get(profile.stage, "No especificado")
    domicilio = (
        f"{profile.domicile.address_road}, {profile.domicile.address_number}, "
        f"{profile.domicile.address_col}, {profile.domicile.address_municip}, "
        f"{profile.domicile.address_city}, {profile.domicile.address_state}, CP {profile.domicile.address_PC}"
    ) if profile.domicile else "N/A"

    user_info = [
        ["Nombre completo:", full_name, "Género:", gender],
        ["Discapacidad:", discapacidad, "Etapa:", stage],
        ["Teléfono:", profile.phone_number or "N/A", "Fecha nacimiento:", birth_date],
        ["Registro:", registration_date, "Tipo de sangre:", profile.blood_type or "N/A"],
        ["CURP:", profile.curp or "N/A", "Alergias:", profile.allergies or "N/A"],
        ["Dirección:", domicilio, "", ""],
    ]

    user_table = Table(user_info, colWidths=[100, 220, 100, 200])
    user_table.setStyle(TableStyle([
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('FONTNAME', (2, 0), (2, -1), 'Helvetica-Bold'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('BACKGROUND', (0, 0), (-1, -1), colors.whitesmoke),
        ('SPAN', (1, 5), (-1, 5)),
    ]))
    elements.append(Spacer(1, 40))
    elements.append(user_table)
    elements.append(Spacer(1, 40))

    # Columnas
    col_widths = [225, 60, 65, 55, 320]

    # Preguntas
    habilidades_laborales = [
        "Asistencia", "Puntualidad", "Higiene y buena presencia", "Avisa cuando llega",
        "Avisa cuando no va a venir", "Avisa cuando termina actividad", "Sigue 1-2 instrucciones",
        "Respeta autoridad", "Es amable con compañeros", "Respeta reglamento", "Trabaja en equipo",
        "No usa celular en el trabajo", "Si no entendió o se le olvidó, pregunta",
        "Ayuda a compañeros de trabajo", "Atención a clientes amable", "Trabaja aún con distracciones",
        "Concluye la tarea asignada", "Revisa que su tarea esté bien hecha",
        "Clasifica y acomoda correctamente", "Ordena correctamente",
        "Recibe nueva mercancía revisando cantidad y calidad",
        "Hace reposición de productos cuando se le pide", "Hace inventario",
        "Conoce y ejecuta reglas de un probador", "Promueve la venta"
    ]

    conducta_adaptativa = [
        "Lectoescritura", "Manejo de dinero", "Comunicación", "Uso de transporte", "Buenos modales",
        "Uso de comunidad", "Guarda distancia correctamente", "Sabe poner límites a agresiones",
        "Sabe negociar", "Toma decisiones importantes en su vida", "Acepta supervisión",
        "Resuelve problemas sencillos", "Maneja correctamente sus emociones", "Se comporta con ética"
    ]

    # --- START OF MODIFIED RESPUESTAS_DICT CREATION ---
    # Fetch responses using the user's specified query
    respuestas_queryset = Respuesta.objects.filter(
        usuario=profile.user.id,
        cuestionario__nombre__iexact="Cuadro de Habilidades",
        cuestionario__activo=True
    ).select_related('pregunta') # Add select_related for efficiency

    # Initialize respuestas_data to store parsed JSON content
    respuestas_data = {}
    for r in respuestas_queryset: # Iterate over the fetched queryset
        normalized_pregunta_text = normalizar(r.pregunta.texto)
        
        parsed_response_content = {}
        if r.respuesta: # Check if the respuesta field has content
            try:
                # Attempt to parse the JSON string
                parsed_response_content = json.loads(r.respuesta)
            except json.JSONDecodeError as e:
                print(f"Error decoding JSON for question '{r.pregunta.texto}': {e}, Raw response: {r.respuesta}")
                # Provide default/empty values if parsing fails
                parsed_response_content = {"resultado": "", "aid_id": None, "aid_text": "Error al procesar apoyo."}
        
        # Store the entire parsed dictionary for each question
        respuestas_data[normalized_pregunta_text] = parsed_response_content
    # --- END OF MODIFIED RESPUESTAS_DICT CREATION ---

    # This map is now for display purposes, mapping the 'resultado' string from JSON
    # to the desired display text for the table columns.
    # The old opciones_map is no longer directly used for determining the 'X' mark,
    # but the display texts might still be useful if you convert 'resultado_raw' for display.
    resultado_display_map = {
        "no_lo_hace": "No lo hace",
        "en_proceso": "En proceso",
        "lo_hace": "Lo hace",
    }

    def generar_tabla(preguntas, titulo):
        data = [[titulo, "No lo hace", "En proceso", "Lo hace", "Apoyos"]]
        # Cache all CHItem objects to avoid repeated queries inside the loop
        all_ch_items = list(CHItem.objects.all()) 

        for pregunta_original in preguntas:
            if not isinstance(pregunta_original, str):
                continue
            
            # Normalize the question text for lookup in respuestas_data
            pregunta_normalized = normalizar(pregunta_original)
            
            # Get the parsed response data for the current question
            # Default to an empty dict if the question has no response data
            response_info = respuestas_data.get(pregunta_normalized, {})
            
            # Extract 'resultado' and 'aid_text' from the parsed JSON data
            resultado_raw = response_info.get('resultado', '')
            aid_text_from_response = response_info.get('aid_text', '')

            # Determine if support text should be shown: if 'resultado' is "no_lo_hace" or "en_proceso"
            mostrar_apoyo = resultado_raw in ["no_lo_hace", "en_proceso"]

            # Set the final aid_text for the 'Apoyos' column
            final_aid_text = ""
            if mostrar_apoyo:
                if aid_text_from_response:
                    final_aid_text = aid_text_from_response
                else:
                    # Fallback to CHItem's aid if aid_text was somehow missing from the response JSON
                    ch_item_obj = next(
                        (item for item in all_ch_items if normalizar(pregunta_original) in normalizar(item.name)),
                        None
                    )
                    final_aid_text = ch_item_obj.aid if ch_item_obj and ch_item_obj.aid else "Sin apoyo registrado"
            
            # Create Paragraph for support text only if there's text to show
            columna_apoyos = Paragraph(final_aid_text, apoyo_style) if final_aid_text else ""

            fila = [
                Paragraph(pregunta_original, left_align),
                "X" if resultado_raw == "no_lo_hace" else "",
                "X" if resultado_raw == "en_proceso" else "",
                "X" if resultado_raw == "lo_hace" else "",
                columna_apoyos,
            ]
            data.append(fila)

        tabla = Table(data, colWidths=col_widths)
        tabla.setStyle(TableStyle([
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('BACKGROUND', (0, 0), (-1, 0), NAVY),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('ALIGN', (1, 1), (-2, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
        ]))
        return tabla

    elements.append(generar_tabla(habilidades_laborales, "HABILIDADES LABORALES"))
    elements.append(PageBreak())
    elements.append(generar_tabla(conducta_adaptativa, "CONDUCTA ADAPTATIVA"))

    elements.append(Spacer(1, 40))
    firma_style = ParagraphStyle(name="FirmaStyle", parent=styles['Normal'], alignment=1, fontSize=10)
    elements.append(Paragraph("______________________________", firma_style))
    elements.append(Paragraph("Responsable 1", firma_style))
    elements.append(Spacer(1, 15))
    elements.append(Paragraph("______________________________", firma_style))
    elements.append(Paragraph("Responsable 2", firma_style))

    doc.build(elements)
    buffer.seek(0)

    return HttpResponse(buffer.getvalue(), content_type="application/pdf", headers={
        "Content-Disposition": f'attachment; filename="cuadro_de_habilidades_{profile.user.id}.pdf"'
    })

def generate_plan_apoyos(p, profile):
    # Stub or actual implementation if available
    p.setTitle("Plan Personalizado de Apoyos")
    p.setFont("Helvetica-Bold", 18)
    p.drawString(200, 700, "PLAN PERSONALIZADO DE APOYOS")
    p.showPage()
    return p